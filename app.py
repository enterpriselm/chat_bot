import time
import streamlit as st
import asyncio
import json
from dotenv import load_dotenv
import google.generativeai as genai
import pandas as pd
import tifffile
from PIL import Image
import PyPDF2
from pptx import Presentation
import cv2
import scipy.io
import ezdxf
import numpy as np
import h5py
from netCDF4 import Dataset
from astropy.io import fits
import avro.datafile
import avro.schema
import trimesh
import open3d as o3d
import os
import io
import yaml # Import the yaml library

# --- Custom Classes (replacing genaitor components) ---
class TaskResult:
    def __init__(self, success: bool, content: any, error: str = None, metadata: dict = None):
        self.success = success
        self.content = content
        self.error = error
        self.metadata = metadata if metadata is not None else {}

    def __str__(self):
        return f"TaskResult(success={self.success}, content={self.content}, error={self.error}, metadata={self.metadata})"


# --- Configuration Loading ---
def load_prompt_config(file_path="config.yaml"):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            config = yaml.safe_load(file)
        return config
    except FileNotFoundError:
        st.error(f"Configuration file '{file_path}' not found. Please create it.")
        st.stop()
    except yaml.YAMLError as e:
        st.error(f"Error parsing YAML file '{file_path}': {e}")
        st.stop()

# Load configuration at the start
PROMPT_CONFIG = load_prompt_config()

# Ensure .env is loaded
load_dotenv(r'.env')

# Access API key
GOOGLE_API_KEY = os.getenv("API_KEY")
if not GOOGLE_API_KEY:
    st.error("`API_KEY` not found in environment variables. Please set it in your `.env` file.")
    st.stop()

genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-1.5-flash')

# --- Task Functions (adapted from classes) ---
def answer_question_task_execute(query: str, context: str, previous_answer: str = None, feedback: str = None) -> TaskResult:
    # st.write("DEBUG: Executando AnswerQuestionTask...") # Removed for cleaner output
    task_config = PROMPT_CONFIG['agents']['answer_question_agent']
    description = task_config['description']
    goal = task_config['goal']
    output_format = "A concise and precise answer to the question." # Defined directly as per user request

    prompt_template = task_config['prompt_template']

    previous_answer_section = f"Previous Answer: {previous_answer}\n" if previous_answer else ""
    feedback_section = f"Feedback for improvement: {feedback}\n" if feedback else ""

    prompt = prompt_template.format(
        description=description,
        goal=goal,
        query=query,
        context=context,
        previous_answer_section=previous_answer_section,
        feedback_section=feedback_section,
        output_format=output_format
    )

    try:
        # st.write("DEBUG: Gerando resposta com LLM (AnswerQuestionTask)...") # Removed for cleaner output
        raw_response = model.generate_content(prompt)
        time.sleep(5)
        response_text = ""
        if raw_response and raw_response.candidates:
            for candidate in raw_response.candidates:
                if candidate.content and candidate.content.parts:
                    for part in candidate.content.parts:
                        if hasattr(part, 'text'):
                            response_text += part.text

        if not response_text:
            raise ValueError("LLM generated no text content.")

        # st.write("DEBUG: Resposta do LLM gerada (AnswerQuestionTask).") # Removed for cleaner output
        return TaskResult(
            success=True,
            content=response_text,
            metadata={"task_type": "answer_question"}
        )
    except Exception as e:
        st.error(f"DEBUG: Erro na AnswerQuestionTask: {e}")
        return TaskResult(success=False, content=None, error=str(e))

def evaluate_answer_task_execute(query: str, context: str, generated_answer: str) -> TaskResult:
    # st.write("DEBUG: Executando EvaluateAnswerTask...") # Removed for cleaner output
    task_config = PROMPT_CONFIG['agents']['evaluate_answer_agent']
    description = task_config['description']
    goal = task_config['goal']
    output_format = task_config['output_format'] # This is the JSON structure

    prompt_template = task_config['prompt_template']
    prompt = prompt_template.format(
        description=description,
        goal=goal,
        query=query,
        context=context,
        generated_answer=generated_answer,
        output_format=output_format
    )

    try:
        # st.write("DEBUG: Gerando avaliação com LLM (EvaluateAnswerTask)...") # Removed for cleaner output
        response = model.generate_content(prompt)
        time.sleep(5)
        response_text = ""
        if response and response.candidates:
            for candidate in response.candidates:
                if candidate.content and candidate.content.parts:
                    for part in candidate.content.parts:
                        if hasattr(part, 'text'):
                            response_text += part.text
        # st.write("DEBUG: Avaliação do LLM gerada (EvaluateAnswerTask).") # Removed for cleaner output
        return TaskResult(
            success=True,
            content=response_text,
            metadata={"task_type": "evaluate_answer"}
        )
    except Exception as e:
        st.error(f"DEBUG: Erro na EvaluateAnswerTask: {e}")
        return TaskResult(success=False, content=None, error=str(e))

def suggest_improvement_task_execute(query: str, context: str, generated_answer: str, evaluation_feedback: str) -> TaskResult:
    # st.write("DEBUG: Executando SuggestImprovementTask...") # Removed for cleaner output
    task_config = PROMPT_CONFIG['agents']['suggest_improvement_agent']
    description = task_config['description']
    goal = task_config['goal']
    output_format = task_config['output_format']

    prompt_template = task_config['prompt_template']
    prompt = prompt_template.format(
        description=description,
        goal=goal,
        query=query,
        context=context,
        generated_answer=generated_answer,
        evaluation_feedback=evaluation_feedback,
        output_format=output_format
    )

    try:
        # st.write("DEBUG: Gerando sugestão de melhoria com LLM (SuggestImprovementTask)...") # Removed for cleaner output
        response = model.generate_content(prompt)
        time.sleep(5)
        response_text = ""
        if response and response.candidates:
            for candidate in response.candidates:
                if candidate.content and candidate.content.parts:
                    for part in candidate.content.parts:
                        if hasattr(part, 'text'):
                            response_text += part.text
        # st.write("DEBUG: Sugestão de melhoria do LLM gerada (SuggestImprovementTask).") # Removed for cleaner output
        return TaskResult(
            success=True,
            content=response_text,
            metadata={"task_type": "suggest_improvement"}
        )
    except Exception as e:
        st.error(f"DEBUG: Erro na SuggestImprovementTask: {e}")
        return TaskResult(success=False, content=None, error=str(e))

# Helper function to extract file content (with the suggested fix)
def extract_file_content(file_path, file_format):
    
    st.write(f"DEBUG: Iniciando extração de conteúdo do arquivo: {file_path} ({file_format})")
    file_content = ""
    # Define the prompt template as a single string
    prompt_template = (
        'You are a highly skilled summarization agent. Your task is to reduce the following {file_format} content to a compact, yet highly informative summary. Focus on preserving all core ideas, characters, events, arguments, and insights necessary to answer detailed questions later. Use a format that balances brevity and informativeness, similar to the Pareto principle: retain the critical 20% that conveys 80% of the value.'
        '\nOutput the summary in a structured format, with sections like:'
        '\n- Title and Metadata (if available)'
        '\n- Key Concepts or Themes'
        '\n- Main Events or Plot Points'
        '\n- Important Characters or Entities'
        '\n- Critical Insights or Takeaways'
        '\nEnsure that no essential information is lost, even if minor details are omitted. The summary should allow deep questioning about the original content without needing to re-read the full text.'
        '\nData: {file_content}'
    )

    try:
        if file_format in ['jpg', 'png']:
            st.write(f"DEBUG: Processando imagem: {file_format}")
            img = Image.open(file_path)
            response = model.generate_content(img)
            time.sleep(5)
            return response.candidates[0].content.parts[0].text
        elif file_format == 'csv':
            st.write("DEBUG: Processando CSV.")
            df = pd.read_csv(file_path)
            file_content = df.to_string()
        elif file_format in ['xls', 'xlsx']:
            st.write("DEBUG: Processando Excel.")
            df = pd.read_excel(file_path)
            file_content = df.to_string()
        elif file_format == 'json':
            st.write("DEBUG: Processando JSON.")
            with open(file_path, 'r', encoding='utf-8') as file: # Added encoding
                data = json.load(file)
            file_content = json.dumps(data, indent=4)
        elif file_format == 'tiff':
            st.write("DEBUG: Processando TIFF.")
            with tifffile.TiffFile(file_path) as tif:
                metadata = tif.pages[0].tags
            file_content = str(metadata)
        elif file_format == 'pdf':
            st.write("DEBUG: Processando PDF.")
            pdf = PyPDF2.PdfReader(file_path)
            for page in pdf.pages:
                file_content += page.extract_text() or "" # Handle empty text
            
        elif file_format == 'ppt' or file_format == 'pptx': # Combine for common PPT formats
            st.write("DEBUG: Processando PPT/PPTX.")
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                file_content += run.text
        elif file_format == 'mp4':
            st.write("DEBUG: Processando MP4.")
            cap = cv2.VideoCapture(file_path)
            frames = []
            while cap.isOpened():
                ret, frame = cap.read()
                if not ret:
                    break
                frames.append(frame)
            cap.release()
            descriptions = []
            frame_sample_rate = max(1, len(frames) // 10)
            for i, frame in enumerate(frames):
                if i % frame_sample_rate == 0:
                    _, buffer = cv2.imencode('.jpg', frame)
                    img = buffer.tobytes()
                    response = model.generate_content(Image.open(io.BytesIO(img)))
                    time.sleep(5)
                    descriptions.append(response.text)
            file_content = "\n".join(descriptions)
        elif file_format == 'mat':
            st.write("DEBUG: Processando MAT.")
            data = scipy.io.loadmat(file_path)
            file_content = str(data)
        elif file_format in ['npy', 'npz']:
            st.write(f"DEBUG: Processando NumPy array: {file_format}.")
            array = np.load(file_path)
            file_content = np.array2string(array)
        elif file_format == 'dxf':
            st.write("DEBUG: Processando DXF.")
            doc = ezdxf.readfile(file_path)
            msp = doc.modelspace()
            data = []
            for entity in msp:
                data.append(str(entity.dxf))
            file_content = "\n".join(data)
        elif file_format == 'hdf5':
            st.write("DEBUG: Processando HDF5.")
            file_content = ""
            with h5py.File(file_path, 'r') as f:
                def print_attrs(name, obj):
                    nonlocal file_content
                    file_content += f"Group or Dataset: {name}\n"
                    for key, val in obj.attrs.items():
                        file_content += f"   Attribute: {key} = {val}\n"
                    if isinstance(obj, h5py.Dataset):
                        file_content += f"   Data: {obj[:]}\n"
                f.visititems(print_attrs)
        elif file_format == 'nc':
            st.write("DEBUG: Processando NetCDF.")
            file_content = ""
            with Dataset(file_path, 'r') as nc:
                for var_name in nc.variables:
                    var = nc.variables[var_name]
                    file_content += f"Variable: {var_name}\n"
                    file_content += f"   Data: {var[:]}\n"
                    for attr_name in var.ncattrs():
                        file_content += f"   Attribute: {attr_name} = {var.getncattr(attr_name)}\n"
        elif file_format == 'fits':
            st.write("DEBUG: Processando FITS.")
            file_content = ""
            with fits.open(file_path) as hdul:
                for hdu in hdul:
                    file_content += f"Extension: {hdu.name}\n"
                    file_content += f"   Header: {hdu.header}\n"
                    if hdu.data is not None:
                        file_content += f"   Data: {hdu.data}\n"
        elif file_format == 'parquet':
            st.write("DEBUG: Processando Parquet.")
            df = pd.read_parquet(file_path)
            file_content = df.to_string()
        elif file_format == 'avro':
            st.write("DEBUG: Processando Avro.")
            file_content = ""
            with open(file_path, 'rb') as fo:
                reader = avro.datafile.Reader(fo, avro.schema.ReadersWriterSchema())
                for record in reader:
                    file_content += str(record) + "\n"
        elif file_format == 'ply':
            st.write("DEBUG: Processando PLY.")
            mesh = trimesh.load(file_path)
            file_content = f"Vertices: {mesh.vertices}\n"
            file_content += f"Faces: {mesh.faces}\n"
        elif file_format == 'pcd':
            st.write("DEBUG: Processando PCD.")
            pcd = o3d.io.read_point_cloud(file_path)
            points = pcd.points
            colors = pcd.colors
            file_content = f"Points: {points}\n"
            file_content += f"Colors: {colors}\n"
        else:
            st.write(f"DEBUG: Formato de arquivo não suportado: .{file_format}")
            return f"Error: Unsupported file format: .{file_format}"

        if file_format not in ['jpg', 'png']:
            st.write("DEBUG: Enviando conteúdo para o modelo Gemini para sumarização.")
            response = model.generate_content(prompt_template.format(file_format=file_format, file_content=file_content))
            time.sleep(5)
            st.write("DEBUG: Sumarização do modelo Gemini concluída.")
            return response.candidates[0].content.parts[0].text
        
    except Exception as e:
        st.error(f"DEBUG: Erro ao processar o arquivo: {e}")
        return f"Error processing file: {e}"
    
    st.write("DEBUG: Extração de conteúdo do arquivo finalizada.")
    return file_content

async def evaluate_until_threshold(query, context, threshold, max_attempts=5):
    st.write("DEBUG: Iniciando processo de refinamento iterativo.")
    feedback = ""
    generated_answer = ""
    st.subheader("Iterative Refinement Process:")
    progress_bar = st.progress(0)

    for attempt in range(1, max_attempts + 1):
        st.write(f"--- **Attempt {attempt}** ---")
        st.write(f"DEBUG: Preparando input para a tentativa {attempt}.")

        # Step 1: Generate Answer
        st.write("DEBUG: Chamando answer_question_task_execute.")
        generator_result = answer_question_task_execute(
            query=query,
            context=context,
            previous_answer=generated_answer,
            feedback=feedback
        )
        if not generator_result.success:
            st.error(f"❌ Error generating answer: {generator_result.error}")
            return None
        generated_answer = generator_result.content
        st.markdown(f"**Generated Answer (Attempt {attempt}):**\n```\n{generated_answer}\n```")

        # Step 2: Evaluate Answer
        st.write("DEBUG: Chamando evaluate_answer_task_execute.")
        evaluator_result = evaluate_answer_task_execute(
            query=query,
            context=context,
            generated_answer=generated_answer
        )
        if not evaluator_result.success:
            st.error(f"❌ Error evaluating answer: {evaluator_result.error}")
            return None
        evaluator_output_raw = evaluator_result.content

        try:
            st.write(f"DEBUG: Tentando parsear saída do avaliador para JSON para a tentativa {attempt}.")
            json_start = evaluator_output_raw.find('{')
            json_end = evaluator_output_raw.rfind('}') + 1
            json_string = evaluator_output_raw[json_start:json_end]
            evaluation = json.loads(json_string)
            st.write(f"DEBUG: Saída do avaliador parseada com sucesso para a tentativa {attempt}.")
        except json.JSONDecodeError as e:
            st.warning(f"⚠️ Failed to parse evaluator output as JSON: {e}. Raw output:\n```json\n{evaluator_output_raw}\n```")
            feedback = f"The previous evaluation output was malformed JSON: '{evaluator_output_raw}'. Please re-evaluate and provide a valid JSON object. Also, provide suggestions for improvement based on the question and context if the answer is not perfect."
            progress_bar.progress(attempt / max_attempts)
            st.write(f"DEBUG: JSON malformado, continuando para a próxima tentativa {attempt}.")
            continue # Continue to next attempt with feedback

        st.markdown(f"**Evaluation:**\n```json\n{json.dumps(evaluation, indent=2)}\n```")

        # Check threshold
        st.write(f"DEBUG: Verificando se os critérios de qualidade foram atendidos para a tentativa {attempt}.")
        hallucinations_check = evaluation.get("hallucinations", "Yes").strip().lower() == threshold["hallucinations"].lower()
        accuracy_check = int(evaluation.get("accuracy", 0)) >= threshold["accuracy"]
        relevance_check = int(evaluation.get("relevance", 0)) >= threshold["relevance"]
        completeness_check = int(evaluation.get("completeness", 0)) >= threshold["completeness"]
        conciseness_check = int(evaluation.get("conciseness", 0)) >= threshold["conciseness"]

        if (accuracy_check and relevance_check and completeness_check and conciseness_check and hallucinations_check):
            st.success("✅ Quality threshold met! Final answer accepted.")
            progress_bar.progress(1.0)
            st.write("DEBUG: Limite de qualidade atingido. Retornando resposta final.")
            return generated_answer
        
        # If not met, generate improvement suggestion
        st.write("DEBUG: Chamando suggest_improvement_task_execute.")
        improvement_result = suggest_improvement_task_execute(
            query=query,
            context=context,
            generated_answer=generated_answer,
            evaluation_feedback=json.dumps(evaluation) # Pass the full evaluation as feedback
        )
        if not improvement_result.success:
            st.error(f"❌ Error suggesting improvement: {improvement_result.error}")
            return None
        feedback = improvement_result.content
        st.info(f"💡 Improvement suggestion for next attempt: {feedback}")
        progress_bar.progress(attempt / max_attempts)
        st.write(f"DEBUG: Limite de qualidade não atingido. Preparando para a próxima tentativa {attempt+1}.")

    st.warning("⚠️ Maximum attempts reached. Returning the best available answer.")
    progress_bar.progress(1.0)
    st.write("DEBUG: Tentativas máximas alcançadas. Retornando a melhor resposta disponível.")
    return generated_answer

# --- Streamlit UI ---

st.set_page_config(page_title="RAG System with Self-Correction", layout="wide")

st.title("📚 RAG System with Self-Correction and Refinement")
st.markdown("""
This application demonstrates a Retrieval-Augmented Generation (RAG) system that
uses an agent orchestrator to extract information from documents,
answer questions, and self-evaluate the quality of the answers, refining them
iteratively until a quality threshold is met.
""")

st.sidebar.header("Document Upload")
uploaded_file = st.sidebar.file_uploader("Choose a file", type=[
    "csv", "xls", "xlsx", "json", "txt", "pdf", "docx", "pptx", "jpg", "png",
    "tiff", "mp4", "mat", "npy", "npz", "dxf", "hdf5", "nc", "fits", "parquet",
    "avro", "ply", "pcd"
], key="file_uploader") # Added a key for better state management

# Initialize document_context and other states
if 'document_context' not in st.session_state:
    st.session_state.document_context = ""
if 'file_processed' not in st.session_state:
    st.session_state.file_processed = False
if 'last_uploaded_file_id' not in st.session_state:
    st.session_state.last_uploaded_file_id = None # Stores the file hash or unique ID

st.write("DEBUG: document_context na sessão:", "Carregado" if st.session_state.document_context else "Vazio")
st.write("DEBUG: file_processed na sessão:", st.session_state.file_processed)
st.write("DEBUG: last_uploaded_file_id na sessão:", st.session_state.last_uploaded_file_id)


# --- Modified logic for document processing ---
if uploaded_file is not None:
    # Use a unique identifier for the file to prevent reprocessing on subsequent reruns
    # A simple way is to use the hash of the file's content or a combination of name and size
    current_file_id = f"{uploaded_file.name}_{uploaded_file.size}" 

    if current_file_id != st.session_state.last_uploaded_file_id:
        st.info("New or different document detected. Processing...")
        st.write("DEBUG: Novo arquivo detectado. Redefinindo estado de processamento.")
        
        # Reset state for new file
        st.session_state.file_processed = False
        st.session_state.document_context = ""
        st.session_state.last_uploaded_file_id = current_file_id # Update the last processed file ID

        temp_file_name = f"temp_uploaded_file_{os.getpid()}_{os.path.splitext(uploaded_file.name)[1]}"
        temp_file_path = os.path.join(temp_file_name)
        with open(temp_file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        file_format = os.path.splitext(uploaded_file.name)[1][1:].lower() # Get format without dot
        with st.spinner("Extracting and summarizing document content..."):
            st.session_state.document_context = extract_file_content(temp_file_path, file_format)
        
        try:
            os.remove(temp_file_path)
            st.write("DEBUG: Arquivo temporário removido com sucesso.")
        except OSError as e:
            st.sidebar.warning(f"Could not remove temporary file: {e}")
            st.write(f"DEBUG: Erro ao remover arquivo temporário: {e}")

        if "Error" in st.session_state.document_context:
            st.error(f"Failed to process document: {st.session_state.document_context}")
            st.session_state.file_processed = False # Mark as not processed if error
        else:
            st.success("Document processed successfully!")
            st.session_state.file_processed = True
            st.write("DEBUG: Documento processado com sucesso e marcado como 'processed'.")
            with st.expander("View Document Summary"):
                st.write(st.session_state.document_context[:1000] + "..." if len(st.session_state.document_context) > 1000 else st.session_state.document_context)
        
        # IMPORTANT: Do NOT call st.rerun() here directly after processing
        # Let Streamlit handle the rerun naturally when the state changes.
        # Calling it here leads to the infinite loop.
    else:
        st.info("Document already processed. Ready to ask questions.")
        st.write("DEBUG: Arquivo já processado (mesmo ID).")
else: # No file uploaded or file_uploader cleared
    if st.session_state.last_uploaded_file_id is not None:
        st.write("DEBUG: Arquivo removido/limpo do uploader. Redefinindo estado.")
        st.session_state.file_processed = False
        st.session_state.document_context = ""
        st.session_state.last_uploaded_file_id = None


st.header("Ask a Question")
user_query = st.text_area("Enter your question about the document:", height=100)

st.header("Evaluation Settings")
col1, col2, col3, col4, col5 = st.columns(5)
with col1:
    acc_thresh = st.number_input("Accuracy (0-10)", min_value=0, max_value=10, value=7)
with col2:
    rel_thresh = st.number_input("Relevance (0-10)", min_value=0, max_value=10, value=7)
with col3:
    comp_thresh = st.number_input("Completeness (0-10)", min_value=0, max_value=10, value=7)
with col4:
    conc_thresh = st.number_input("Conciseness (0-10)", min_value=0, max_value=10, value=7)
with col5:
    hall_thresh = st.selectbox("Hallucinations", ["No", "Yes"], index=0)

evaluation_threshold = {
    "accuracy": acc_thresh,
    "relevance": rel_thresh,
    "completeness": comp_thresh,
    "conciseness": conc_thresh,
    "hallucinations": hall_thresh
}

max_attempts = st.slider("Maximum Refinement Attempts", min_value=1, max_value=10, value=5)


# This button now only triggers the question answering, not document processing
if st.button("Generate Answer and Evaluate", disabled=not st.session_state.file_processed):
    st.write("DEBUG: Botão 'Generate Answer and Evaluate' clicado.")
    if not st.session_state.file_processed or not st.session_state.document_context:
        st.warning("Please upload and process a document first.")
        st.write("DEBUG: Documento não processado ou contexto vazio.")
    elif not user_query:
        st.warning("Please enter a question.")
        st.write("DEBUG: Nenhuma pergunta inserida.")
    else:
        st.info("Starting RAG refinement process for your question...")
        st.write("DEBUG: Iniciando processo RAG para a pergunta.")

        final_answer_placeholder = st.empty()
        
        st.write("DEBUG: Iniciando a avaliação até o limite.")
        final_answer = asyncio.run(
            evaluate_until_threshold(
                query=user_query,
                context=st.session_state.document_context,
                threshold=evaluation_threshold,
                max_attempts=max_attempts
            )
        )
        st.write("DEBUG: Avaliação até o limite concluída.")
        
        if final_answer:
            final_answer_placeholder.markdown("---")
            final_answer_placeholder.subheader("🎉 Final Answer Accepted:")
            final_answer_placeholder.success(final_answer)
            st.write("DEBUG: Resposta final aceita e exibida.")
        else:
            final_answer_placeholder.error("Could not generate a satisfactory answer after the maximum attempts.")
            st.write("DEBUG: Não foi possível gerar uma resposta satisfatória.")
